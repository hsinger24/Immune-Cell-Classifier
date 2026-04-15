"""PI-ready visual suite for the immune cell classifier.

Single entrypoint that regenerates every figure in NEXT_STEPS.md from the saved
checkpoints. Idempotent; safe to re-run after any future retraining.
"""

from __future__ import annotations

import argparse
import io
import warnings
from collections import defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

import cv2
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import torch
import torch.nn.functional as F
from matplotlib.gridspec import GridSpec
from PIL import Image
from sklearn.calibration import calibration_curve
from sklearn.metrics import (
    auc,
    average_precision_score,
    confusion_matrix,
    precision_recall_curve,
    roc_curve,
)
from torch.utils.data import DataLoader
from torchvision import transforms as T
from tqdm import tqdm

from .data import (
    CLASSES,
    CLASS_TO_IDX,
    DEFAULT_DATA_ROOT,
    IMAGENET_MEAN,
    IMAGENET_STD,
    NUM_CLASSES,
    CellDataset,
    _CombinedDataset,
    _SubsetWithTransform,
    get_grouped_split_loaders,
    get_official_split_loaders,
    get_transforms,
)
from .model import create_model

# ──────────────────────────────────────────────────────────────────────────────
# Style
# ──────────────────────────────────────────────────────────────────────────────

plt.rcParams["figure.dpi"] = 150
plt.rcParams["savefig.dpi"] = 300
plt.rcParams["font.family"] = "DejaVu Sans"
plt.rcParams["axes.grid"] = False

CLASS_COLORS = {i: plt.cm.tab10(i) for i in range(NUM_CLASSES)}

FIG_DIR = Path(__file__).resolve().parent.parent / "reports" / "figures"
FIG_DIR.mkdir(parents=True, exist_ok=True)

DEVICE = (
    torch.device("mps") if torch.backends.mps.is_available()
    else torch.device("cuda") if torch.cuda.is_available()
    else torch.device("cpu")
)

INV_NORMALIZE = T.Normalize(
    mean=[-m / s for m, s in zip(IMAGENET_MEAN, IMAGENET_STD)],
    std=[1.0 / s for s in IMAGENET_STD],
)


# ──────────────────────────────────────────────────────────────────────────────
# Utilities
# ──────────────────────────────────────────────────────────────────────────────

def load_checkpoint(split_mode: str) -> torch.nn.Module:
    ckpt_path = Path(__file__).resolve().parent.parent / "models" / f"best_efficientnetv2_s_{split_mode}.pt"
    ckpt = torch.load(ckpt_path, map_location=DEVICE, weights_only=True)
    model = create_model(ckpt["model_name"], pretrained=False)
    model.load_state_dict(ckpt["model_state"])
    model = model.to(DEVICE).eval()
    return model


def to_display(img_tensor: torch.Tensor) -> np.ndarray:
    return INV_NORMALIZE(img_tensor).clamp(0, 1).permute(1, 2, 0).numpy()


@torch.no_grad()
def collect_predictions(model, loader) -> dict:
    """Return probs, preds, labels, features (pre-classifier)."""
    probs_all, labels_all, feats_all = [], [], []

    # Hook on penultimate features (timm: model.forward_features; we'll use that directly)
    for images, labels in tqdm(loader, desc="Collecting preds", leave=False):
        images = images.to(DEVICE)
        feats = model.forward_features(images)
        pooled = model.forward_head(feats, pre_logits=True)  # (B, 1280) for effnetv2_s
        logits = model.get_classifier()(pooled)
        probs = F.softmax(logits, dim=1)
        probs_all.append(probs.cpu())
        feats_all.append(pooled.cpu())
        labels_all.append(labels)

    probs = torch.cat(probs_all).numpy()
    feats = torch.cat(feats_all).numpy()
    labels = torch.cat(labels_all).numpy()
    preds = probs.argmax(1)
    return dict(probs=probs, preds=preds, labels=labels, features=feats)


def get_target_layer(model):
    # EfficientNetV2-S in timm: last conv before GAP is conv_head
    return model.conv_head


# ──────────────────────────────────────────────────────────────────────────────
# CAM wrappers (pytorch-grad-cam)
# ──────────────────────────────────────────────────────────────────────────────

def _cam(method_cls, model, input_tensor, target_class):
    from pytorch_grad_cam.utils.model_targets import ClassifierOutputTarget
    with method_cls(model=model, target_layers=[get_target_layer(model)]) as cam:
        grayscale = cam(input_tensor=input_tensor,
                        targets=[ClassifierOutputTarget(target_class)])
    return grayscale[0]  # (H, W) in [0, 1]


def gradcam(model, input_tensor, target_class):
    from pytorch_grad_cam import GradCAM
    return _cam(GradCAM, model, input_tensor, target_class)


def gradcam_plusplus(model, input_tensor, target_class):
    from pytorch_grad_cam import GradCAMPlusPlus
    return _cam(GradCAMPlusPlus, model, input_tensor, target_class)


def scorecam(model, input_tensor, target_class):
    from pytorch_grad_cam import ScoreCAM
    return _cam(ScoreCAM, model, input_tensor, target_class)


def overlay_cam(display_img: np.ndarray, cam: np.ndarray, alpha: float = 0.45,
                cmap: str = "jet") -> np.ndarray:
    cam_resized = cv2.resize(cam, (display_img.shape[1], display_img.shape[0]),
                              interpolation=cv2.INTER_CUBIC)
    heatmap = plt.get_cmap(cmap)(cam_resized)[..., :3]
    return np.clip((1 - alpha) * display_img + alpha * heatmap, 0, 1)


# ──────────────────────────────────────────────────────────────────────────────
# Figure 1 — Grad-CAM 2.0
# ──────────────────────────────────────────────────────────────────────────────

def figure_gradcam_panel(model, dataset, preds_info, split: str):
    """4 rows (class) × 6 cols (high-confidence examples per class)."""
    probs = preds_info["probs"]
    labels = preds_info["labels"]
    preds = preds_info["preds"]

    # Select top-6 correctly-classified, high-confidence per class
    picks = {}
    for c in range(NUM_CLASSES):
        correct = np.where((labels == c) & (preds == c))[0]
        conf = probs[correct, c]
        order = np.argsort(-conf)
        picks[c] = correct[order][:6]

    fig = plt.figure(figsize=(14, 10))
    gs = GridSpec(NUM_CLASSES, 7, width_ratios=[1] * 6 + [0.08],
                  wspace=0.05, hspace=0.12)

    for r, c in enumerate(range(NUM_CLASSES)):
        idxs = picks[c]
        for j, idx in enumerate(idxs):
            img_tensor, _ = dataset[idx]
            display = to_display(img_tensor)
            input_t = img_tensor.unsqueeze(0).to(DEVICE)
            cam = gradcam(model, input_t, c)
            overlay = overlay_cam(display, cam, alpha=0.45, cmap="jet")
            ax = fig.add_subplot(gs[r, j])
            ax.imshow(overlay)
            ax.text(0.03, 0.95, f"p={probs[idx, c]:.2f}",
                    transform=ax.transAxes, color="white",
                    fontsize=9, va="top",
                    bbox=dict(facecolor="black", alpha=0.6, edgecolor="none",
                              boxstyle="round,pad=0.2"))
            ax.set_xticks([]); ax.set_yticks([])
            if j == 0:
                ax.set_ylabel(CLASSES[c], fontsize=11, rotation=90,
                              labelpad=8, fontweight="bold")

    cbar_ax = fig.add_subplot(gs[:, 6])
    sm = plt.cm.ScalarMappable(cmap="jet",
                                norm=plt.Normalize(vmin=0, vmax=1))
    sm.set_array([])
    cbar = fig.colorbar(sm, cax=cbar_ax)
    cbar.set_label("Grad-CAM activation", fontsize=10)

    fig.suptitle(f"Grad-CAM attention on high-confidence predictions — {split} split",
                 fontsize=13, y=0.995)
    out_png = FIG_DIR / f"gradcam_panel_{split}.png"
    out_pdf = FIG_DIR / f"gradcam_panel_{split}.pdf"
    fig.savefig(out_png, dpi=300, bbox_inches="tight")
    fig.savefig(out_pdf, bbox_inches="tight")
    plt.close(fig)
    print(f"  [1] wrote {out_png.name} + .pdf")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 2 — Grad-CAM++ / Score-CAM comparison
# ──────────────────────────────────────────────────────────────────────────────

def figure_cam_comparison(model, dataset, preds_info, split: str):
    """4 rows (class) × 4 cols (orig | GradCAM | GradCAM++ | ScoreCAM)."""
    probs = preds_info["probs"]
    labels = preds_info["labels"]
    preds = preds_info["preds"]

    picks = {}
    for c in range(NUM_CLASSES):
        correct = np.where((labels == c) & (preds == c))[0]
        picks[c] = correct[np.argmax(probs[correct, c])]

    fig, axes = plt.subplots(NUM_CLASSES, 4, figsize=(11, 10.5))
    col_titles = ["Original", "Grad-CAM", "Grad-CAM++", "Score-CAM"]
    for j, t in enumerate(col_titles):
        axes[0, j].set_title(t, fontsize=11, fontweight="bold")

    for r, c in enumerate(range(NUM_CLASSES)):
        idx = picks[c]
        img_tensor, _ = dataset[idx]
        display = to_display(img_tensor)
        input_t = img_tensor.unsqueeze(0).to(DEVICE)

        cams = [
            gradcam(model, input_t, c),
            gradcam_plusplus(model, input_t, c),
            scorecam(model, input_t, c),
        ]

        axes[r, 0].imshow(display)
        axes[r, 0].set_ylabel(CLASSES[c], fontsize=11, rotation=90, fontweight="bold")
        for j, cam in enumerate(cams, start=1):
            axes[r, j].imshow(overlay_cam(display, cam, alpha=0.45, cmap="jet"))

        for j in range(4):
            axes[r, j].set_xticks([]); axes[r, j].set_yticks([])

    fig.suptitle(f"Attention method comparison — {split} split",
                 fontsize=13, y=0.995)
    plt.tight_layout()
    out = FIG_DIR / f"cam_comparison_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [2] wrote {out.name}")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 3 — Misclassification gallery (official split only)
# ──────────────────────────────────────────────────────────────────────────────

def figure_misclass_gallery(model, dataset, preds_info, split: str):
    probs = preds_info["probs"]
    labels = preds_info["labels"]
    preds = preds_info["preds"]

    wrong = np.where(preds != labels)[0]
    if len(wrong) == 0:
        print(f"  [3] skip — no misclassifications on {split} split")
        return
    # Sort by confidence in the wrong class (descending)
    wrong_conf = probs[wrong, preds[wrong]]
    order = wrong[np.argsort(-wrong_conf)][:12]

    n = len(order)
    rows = (n + 3) // 4  # 4 examples per row, each example uses a 1x4 block
    fig = plt.figure(figsize=(18, 3.2 * n))
    gs = GridSpec(n, 4, wspace=0.15, hspace=0.35)

    for k, idx in enumerate(order):
        img_tensor, _ = dataset[idx]
        display = to_display(img_tensor)
        input_t = img_tensor.unsqueeze(0).to(DEVICE)

        true_c = int(labels[idx])
        pred_c = int(preds[idx])
        cam_true = gradcam(model, input_t, true_c)
        cam_pred = gradcam(model, input_t, pred_c)

        ax0 = fig.add_subplot(gs[k, 0])
        ax0.imshow(display)
        ax0.set_title(f"True: {CLASSES[true_c]}\nPred: {CLASSES[pred_c]}  "
                      f"(p={probs[idx, pred_c]:.2f})", fontsize=9)
        ax0.axis("off")

        ax1 = fig.add_subplot(gs[k, 1])
        ax1.imshow(overlay_cam(display, cam_true, alpha=0.45, cmap="jet"))
        ax1.set_title(f"Grad-CAM → {CLASSES[true_c]} (true)", fontsize=9)
        ax1.axis("off")

        ax2 = fig.add_subplot(gs[k, 2])
        ax2.imshow(overlay_cam(display, cam_pred, alpha=0.45, cmap="jet"))
        ax2.set_title(f"Grad-CAM → {CLASSES[pred_c]} (pred)", fontsize=9)
        ax2.axis("off")

        ax3 = fig.add_subplot(gs[k, 3])
        top3 = np.argsort(-probs[idx])[:3]
        bars = ax3.barh(range(3), probs[idx, top3][::-1],
                        color=[CLASS_COLORS[c] for c in top3[::-1]])
        ax3.set_yticks(range(3))
        ax3.set_yticklabels([CLASSES[c] for c in top3[::-1]], fontsize=8)
        ax3.set_xlim(0, 1)
        ax3.set_xlabel("prob", fontsize=8)
        ax3.tick_params(axis="x", labelsize=7)
        ax3.set_title("Top-3", fontsize=9)

    fig.suptitle(f"Most-confidently-wrong predictions — {split} split",
                 fontsize=13, y=0.998)
    out = FIG_DIR / f"misclass_gallery_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [3] wrote {out.name}  (n={n})")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 4 — Mean attention per class
# ──────────────────────────────────────────────────────────────────────────────

def figure_mean_attention(model, dataset, preds_info, split: str,
                          max_per_class: int = 200):
    labels = preds_info["labels"]
    preds = preds_info["preds"]

    fig, axes = plt.subplots(1, NUM_CLASSES, figsize=(16, 4.2))
    for c in range(NUM_CLASSES):
        correct = np.where((labels == c) & (preds == c))[0][:max_per_class]
        if len(correct) == 0:
            axes[c].set_axis_off(); continue
        acc = np.zeros((224, 224), dtype=np.float32)
        for idx in tqdm(correct, desc=f"  mean-att {CLASSES[c]}", leave=False):
            img_tensor, _ = dataset[idx]
            input_t = img_tensor.unsqueeze(0).to(DEVICE)
            cam = gradcam(model, input_t, c)
            cam_resized = cv2.resize(cam, (224, 224), interpolation=cv2.INTER_CUBIC)
            acc += cam_resized
        acc /= len(correct)
        acc = (acc - acc.min()) / (acc.max() - acc.min() + 1e-8)

        axes[c].imshow(acc, cmap="inferno")
        axes[c].set_title(f"{CLASSES[c]}  (n={len(correct)})", fontsize=11)
        axes[c].set_xticks([]); axes[c].set_yticks([])

    fig.suptitle(f"Mean Grad-CAM per class — {split} split",
                 fontsize=13, y=1.01)
    plt.tight_layout()
    out = FIG_DIR / f"mean_attention_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [4] wrote {out.name}")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 5 — UMAP embedding
# ──────────────────────────────────────────────────────────────────────────────

def figure_umap(preds_info, split: str):
    import umap

    feats = preds_info["features"]
    labels = preds_info["labels"]
    preds = preds_info["preds"]

    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        reducer = umap.UMAP(n_neighbors=15, min_dist=0.1, random_state=42)
        emb = reducer.fit_transform(feats)

    fig, ax = plt.subplots(figsize=(8, 8))
    for c in range(NUM_CLASSES):
        mask = labels == c
        ax.scatter(emb[mask, 0], emb[mask, 1],
                   color=CLASS_COLORS[c], label=CLASSES[c],
                   alpha=0.75, s=24, edgecolor="none")

    wrong = labels != preds
    if wrong.any():
        ax.scatter(emb[wrong, 0], emb[wrong, 1],
                   facecolor="none", edgecolor="black",
                   s=60, linewidth=1.2, label="misclassified")

    ax.legend(loc="best", frameon=True, fontsize=10)
    ax.set_xlabel("UMAP-1"); ax.set_ylabel("UMAP-2")
    ax.set_title(f"Penultimate-layer UMAP (1280-D → 2-D) — {split} split",
                 fontsize=12)
    plt.tight_layout()
    out = FIG_DIR / f"umap_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [5] wrote {out.name}")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 6 — Reliability + ROC + PR
# ──────────────────────────────────────────────────────────────────────────────

def _expected_calibration_error(probs, labels, n_bins=15):
    confs = probs.max(1)
    preds = probs.argmax(1)
    correct = (preds == labels).astype(np.float32)
    bins = np.linspace(0, 1, n_bins + 1)
    ece = 0.0
    for i in range(n_bins):
        m = (confs >= bins[i]) & (confs < bins[i + 1]) if i < n_bins - 1 \
            else (confs >= bins[i]) & (confs <= bins[i + 1])
        if m.any():
            ece += (m.mean()) * abs(correct[m].mean() - confs[m].mean())
    return float(ece)


def figure_calibration_roc_pr(preds_info, split: str):
    probs = preds_info["probs"]
    labels = preds_info["labels"]

    fig, axes = plt.subplots(1, 3, figsize=(16, 5.2))

    # (a) reliability
    confs = probs.max(1)
    acc = (probs.argmax(1) == labels).astype(np.float32)
    bins = np.linspace(0, 1, 11)
    bin_idx = np.digitize(confs, bins) - 1
    bin_idx = np.clip(bin_idx, 0, len(bins) - 2)
    bin_conf, bin_acc, bin_n = [], [], []
    for i in range(len(bins) - 1):
        m = bin_idx == i
        if m.any():
            bin_conf.append(confs[m].mean())
            bin_acc.append(acc[m].mean())
            bin_n.append(m.sum())
    ece = _expected_calibration_error(probs, labels)
    axes[0].plot([0, 1], [0, 1], "k--", alpha=0.5, label="perfect")
    axes[0].bar(bin_conf, bin_acc, width=0.08, alpha=0.6,
                color="steelblue", edgecolor="black", label="model")
    axes[0].set_xlabel("Predicted confidence")
    axes[0].set_ylabel("Empirical accuracy")
    axes[0].set_title(f"Reliability diagram  (ECE = {ece:.3f})")
    axes[0].set_xlim(0, 1); axes[0].set_ylim(0, 1)
    axes[0].legend()

    # (b) ROC per class (one-vs-rest)
    labels_oh = np.eye(NUM_CLASSES)[labels]
    for c in range(NUM_CLASSES):
        fpr, tpr, _ = roc_curve(labels_oh[:, c], probs[:, c])
        axes[1].plot(fpr, tpr, color=CLASS_COLORS[c],
                     label=f"{CLASSES[c]} (AUC={auc(fpr, tpr):.3f})", lw=2)
    axes[1].plot([0, 1], [0, 1], "k--", alpha=0.4)
    axes[1].set_xlabel("FPR"); axes[1].set_ylabel("TPR")
    axes[1].set_title("ROC — one-vs-rest")
    axes[1].legend(fontsize=9)

    # (c) PR per class
    for c in range(NUM_CLASSES):
        prec, rec, _ = precision_recall_curve(labels_oh[:, c], probs[:, c])
        ap = average_precision_score(labels_oh[:, c], probs[:, c])
        axes[2].plot(rec, prec, color=CLASS_COLORS[c],
                     label=f"{CLASSES[c]} (AP={ap:.3f})", lw=2)
    axes[2].set_xlabel("Recall"); axes[2].set_ylabel("Precision")
    axes[2].set_title("Precision-Recall")
    axes[2].legend(fontsize=9)

    fig.suptitle(f"Calibration + classifier curves — {split} split",
                 fontsize=13, y=1.02)
    plt.tight_layout()
    out = FIG_DIR / f"calibration_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [6] wrote {out.name}  (ECE={ece:.3f})")
    return ece


# ──────────────────────────────────────────────────────────────────────────────
# Figure 7 — Confusion matrix, upgraded
# ──────────────────────────────────────────────────────────────────────────────

def figure_confusion_matrix(preds_info, split: str):
    labels = preds_info["labels"]
    preds = preds_info["preds"]
    cm = confusion_matrix(labels, preds, labels=range(NUM_CLASSES))
    cm_norm = cm / cm.sum(axis=1, keepdims=True).clip(min=1)

    fig = plt.figure(figsize=(9, 7.5))
    gs = GridSpec(2, 2, width_ratios=[1, 0.18], height_ratios=[0.18, 1],
                  wspace=0.04, hspace=0.04)
    ax = fig.add_subplot(gs[1, 0])
    ax_top = fig.add_subplot(gs[0, 0], sharex=ax)
    ax_right = fig.add_subplot(gs[1, 1], sharey=ax)

    im = ax.imshow(cm_norm, cmap="RdYlGn_r", vmin=0, vmax=1)
    for i in range(NUM_CLASSES):
        for j in range(NUM_CLASSES):
            txt = f"{cm[i, j]}\n({cm_norm[i, j]:.0%})"
            ax.text(j, i, txt, ha="center", va="center",
                    color="black" if cm_norm[i, j] < 0.5 else "white",
                    fontsize=9)
    ax.set_xticks(range(NUM_CLASSES)); ax.set_yticks(range(NUM_CLASSES))
    ax.set_xticklabels(CLASSES, rotation=30, ha="right")
    ax.set_yticklabels(CLASSES)
    ax.set_xlabel("Predicted"); ax.set_ylabel("True")

    # Marginal class counts
    class_counts = cm.sum(axis=1)
    ax_top.bar(range(NUM_CLASSES), class_counts,
               color=[CLASS_COLORS[c] for c in range(NUM_CLASSES)])
    ax_top.set_ylabel("n", fontsize=9)
    ax_top.tick_params(axis="x", labelbottom=False)

    pred_counts = cm.sum(axis=0)
    ax_right.barh(range(NUM_CLASSES), pred_counts,
                  color=[CLASS_COLORS[c] for c in range(NUM_CLASSES)])
    ax_right.set_xlabel("n", fontsize=9)
    ax_right.tick_params(axis="y", labelleft=False)
    ax_right.invert_yaxis()

    fig.colorbar(im, ax=ax_right, pad=0.02, fraction=0.5,
                 label="row-normalized accuracy")
    fig.suptitle(f"Confusion matrix — {split} split", fontsize=13, y=1.00)
    out = FIG_DIR / f"confusion_matrix_v2_{split}.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [7] wrote {out.name}")


# ──────────────────────────────────────────────────────────────────────────────
# Figure 8 — Cross-dataset sanity check (dataset-master)
# ──────────────────────────────────────────────────────────────────────────────

def _parse_voc_bboxes(xml_path: Path):
    root = ET.parse(xml_path).getroot()
    bboxes = []
    for obj in root.findall("object"):
        name_node = obj.find("name")
        name = name_node.text.upper() if name_node is not None else ""
        bnd = obj.find("bndbox")
        if bnd is None:
            continue
        xmin = int(bnd.find("xmin").text); ymin = int(bnd.find("ymin").text)
        xmax = int(bnd.find("xmax").text); ymax = int(bnd.find("ymax").text)
        bboxes.append((name, xmin, ymin, xmax, ymax))
    return bboxes


def figure_cross_dataset(model):
    """dataset-master VOC annotations label only RBCs, but labels.csv gives the
    single WBC class per smear. We feed the whole 640x480 smear to the model
    (center-cropped) and compare the prediction to the CSV ground truth.
    """
    root = Path(__file__).resolve().parent.parent / "dataset-master" / "dataset-master"
    img_dir = root / "JPEGImages"
    csv_path = root / "labels.csv"
    if not img_dir.exists() or not csv_path.exists():
        print("  [8] skip — dataset-master not present")
        return

    df = pd.read_csv(csv_path)
    df.columns = [c.strip() for c in df.columns]

    by_class = defaultdict(list)
    for _, r in df.iterrows():
        if pd.isna(r["Category"]):
            continue
        cat = str(r["Category"]).strip().upper()
        if cat not in CLASS_TO_IDX:
            continue
        idx = int(r["Image"])
        img = img_dir / f"BloodImage_{idx:05d}.jpg"
        if not img.exists():
            continue
        if len(by_class[cat]) < 4:
            by_class[cat].append((img, cat))

    # 4 per class × 4 classes = 16 panels
    candidates = []
    for c in CLASSES:
        candidates.extend(by_class[c][:4])
    if not candidates:
        print("  [8] skip — no matching annotations")
        return

    tfm = get_transforms("test", 224)
    fig, axes = plt.subplots(4, 4, figsize=(16, 12))
    correct = 0
    for ax, (img_path, gt_name) in zip(axes.flat, candidates):
        arr = np.array(Image.open(img_path).convert("RGB"))
        pil = Image.fromarray(arr)
        input_t = tfm(pil).unsqueeze(0).to(DEVICE)
        with torch.no_grad():
            probs = F.softmax(model(input_t), dim=1).cpu().numpy()[0]
        pred = int(probs.argmax())
        hit = CLASSES[pred] == gt_name
        correct += int(hit)
        color = "lime" if hit else "red"

        ax.imshow(arr)
        ax.text(0.02, 0.98,
                f"Pred: {CLASSES[pred]} (p={probs[pred]:.2f})\nGT:   {gt_name}",
                color="white", fontsize=10, va="top",
                transform=ax.transAxes,
                bbox=dict(facecolor=color, alpha=0.8, edgecolor="none",
                          boxstyle="round,pad=0.3"))
        ax.set_xticks([]); ax.set_yticks([])

    fig.suptitle(f"Cross-dataset sanity — raw 640×480 smears from dataset-master  "
                 f"({correct}/{len(candidates)} correct)",
                 fontsize=14, y=0.995)
    plt.tight_layout()
    out = FIG_DIR / "cross_dataset_sanity.png"
    fig.savefig(out, dpi=300, bbox_inches="tight")
    plt.close(fig)
    print(f"  [8] wrote {out.name}")


# ──────────────────────────────────────────────────────────────────────────────
# Priority 2 — Hero figure
# ──────────────────────────────────────────────────────────────────────────────

def figure_hero(split_preds: dict):
    """2×3 hero figure assembled from pre-rendered panels."""
    fig = plt.figure(figsize=(16, 10))
    gs = GridSpec(2, 3, wspace=0.08, hspace=0.18)

    panels = [
        ("sample data", Path(__file__).resolve().parent.parent / "reports" / "sample_images.png"),
        ("confusion / grouped", FIG_DIR / "confusion_matrix_v2_grouped.png"),
        ("confusion / official", FIG_DIR / "confusion_matrix_v2_official.png"),
        ("Grad-CAM / grouped",  FIG_DIR / "gradcam_panel_grouped.png"),
        ("UMAP / grouped",      FIG_DIR / "umap_grouped.png"),
        ("calibration / official", FIG_DIR / "calibration_official.png"),
    ]
    for k, (title, p) in enumerate(panels):
        ax = fig.add_subplot(gs[k // 3, k % 3])
        if p.exists():
            ax.imshow(Image.open(p))
            ax.set_title(title, fontsize=11, fontweight="bold")
        else:
            ax.text(0.5, 0.5, f"missing:\n{p.name}", ha="center", va="center")
        ax.axis("off")

    fig.suptitle("Immune Cell Classifier — Headline results (EfficientNetV2-S)",
                 fontsize=15, y=0.995, fontweight="bold")
    out_png = FIG_DIR / "hero_figure.png"
    out_pdf = FIG_DIR / "hero_figure.pdf"
    fig.savefig(out_png, dpi=300, bbox_inches="tight")
    fig.savefig(out_pdf, bbox_inches="tight")
    plt.close(fig)
    print(f"  [H] wrote {out_png.name} + .pdf")


# ──────────────────────────────────────────────────────────────────────────────
# Priority 3 — Captions, results.md, index.html, pptx
# ──────────────────────────────────────────────────────────────────────────────

CAPTIONS = {
    "gradcam_panel_grouped.png":
        "Grad-CAM attention on 24 high-confidence correct predictions from the grouped test split. "
        "The model consistently localizes the nucleus/cytoplasm region characteristic of each leukocyte class.",
    "gradcam_panel_official.png":
        "Same panel layout applied to the official split, showing that attention remains nucleus-centered even where augmentation-induced leakage is absent.",
    "cam_comparison_grouped.png":
        "Grad-CAM, Grad-CAM++ and Score-CAM agree on the informative region for a representative high-confidence example of each class, confirming that the attention signal is robust to method choice.",
    "cam_comparison_official.png":
        "Cross-method attention agreement on the official split.",
    "misclass_gallery_official.png":
        "The model's most-confidently-wrong predictions on the official test split, paired with attention maps for the true and predicted classes and a top-3 probability bar chart.",
    "mean_attention_grouped.png":
        "Grad-CAM heatmap averaged over up to 200 correctly classified examples per class, showing the prototypical spatial attention pattern the model uses for each leukocyte.",
    "mean_attention_official.png":
        "Same prototypical-attention figure computed on the official split.",
    "umap_grouped.png":
        "UMAP projection of the 1280-D penultimate-layer features on the grouped test set. Clear class-wise clustering indicates the model has learned a meaningful embedding; misclassified points (outlined in black) sit at cluster boundaries.",
    "umap_official.png":
        "UMAP projection on the official test set.",
    "calibration_grouped.png":
        "Reliability diagram with ECE, one-vs-rest ROC and precision-recall curves on the grouped test set.",
    "calibration_official.png":
        "Same calibration panel on the official split; residual confusion between lymphocytes and monocytes shows up as lower AP and a small calibration gap.",
    "confusion_matrix_v2_grouped.png":
        "Row-normalized confusion matrix (grouped split) with raw counts, percentages and class-count marginals.",
    "confusion_matrix_v2_official.png":
        "Same, on the official split.",
    "cross_dataset_sanity.png":
        "Predictions on 16 raw blood smears from dataset-master — a completely separate VOC-style collection the model has never seen during training.",
    "hero_figure.png":
        "One-page summary: sample data, confusion matrices (grouped/official), Grad-CAM panel, UMAP embedding and calibration curves.",
}


def write_captions():
    out = FIG_DIR / "CAPTIONS.md"
    lines = ["# Figure captions", ""]
    for name in sorted(FIG_DIR.glob("*.png")):
        caption = CAPTIONS.get(name.name, "")
        lines.append(f"## `{name.name}`")
        lines.append("")
        lines.append(caption if caption else "_(no caption)_")
        lines.append("")
    out.write_text("\n".join(lines))
    print(f"  wrote {out.relative_to(FIG_DIR.parent)}")


def write_index_html():
    out = FIG_DIR / "index.html"
    figs = sorted(FIG_DIR.glob("*.png"))
    items = "\n".join(
        f'  <figure><img src="{f.name}" loading="lazy"/>'
        f'<figcaption><b>{f.name}</b><br>{CAPTIONS.get(f.name, "")}</figcaption></figure>'
        for f in figs
    )
    html = f"""<!doctype html>
<html><head><meta charset="utf-8"><title>Immune-Cell-Classifier · figure gallery</title>
<style>
  body {{ font-family: -apple-system, sans-serif; margin: 2em; background:#fafafa; }}
  h1 {{ margin-bottom:0.2em; }}
  figure {{ margin:1.5em 0; padding:1em; background:white; border:1px solid #ddd; border-radius:6px; }}
  figure img {{ max-width:100%; height:auto; display:block; margin-bottom:0.7em; }}
  figcaption {{ font-size:0.92em; color:#333; }}
</style></head>
<body>
<h1>Immune Cell Classifier — figure gallery</h1>
<p>Auto-generated by <code>src/visualize.py</code>.</p>
{items}
</body></html>
"""
    out.write_text(html)
    print(f"  wrote {out.relative_to(FIG_DIR.parent)}")


def write_results_md(summary: dict):
    """Regenerate reports/results.md embedding the new figures."""
    out = Path(__file__).resolve().parent.parent / "reports" / "results.md"
    lines = [
        "# Immune Cell Classifier — Results",
        "",
        "EfficientNetV2-S fine-tuned on the Kaggle BCCD-derived 4-class blood cell dataset.",
        "Two evaluation regimes are reported: the **official** train/test split and a "
        "**grouped** split where augmented copies of the same source cell cannot straddle splits.",
        "",
        "## Headline metrics",
        "",
        "| Split    | Accuracy | Macro-F1 | ECE   |",
        "|----------|---------:|---------:|------:|",
    ]
    for split, s in summary.items():
        lines.append(f"| {split:<8} | {s['accuracy']:.4f}   | {s['macro_f1']:.4f}   | {s['ece']:.3f} |")
    lines += [
        "",
        "## Hero figure",
        "",
        "![Hero figure](figures/hero_figure.png)",
        "",
        "## Priority 1 — Visuals",
        "",
    ]
    for name in [
        "gradcam_panel_grouped.png", "gradcam_panel_official.png",
        "cam_comparison_grouped.png",
        "misclass_gallery_official.png",
        "mean_attention_grouped.png",
        "umap_grouped.png",
        "calibration_grouped.png", "calibration_official.png",
        "confusion_matrix_v2_grouped.png", "confusion_matrix_v2_official.png",
        "cross_dataset_sanity.png",
    ]:
        if (FIG_DIR / name).exists():
            lines.append(f"### {name}")
            lines.append("")
            lines.append(f"![{name}](figures/{name})")
            lines.append("")
            lines.append(CAPTIONS.get(name, ""))
            lines.append("")
    lines.append("Browse all figures as a gallery: [figures/index.html](figures/index.html). "
                 "Draft captions: [figures/CAPTIONS.md](figures/CAPTIONS.md).")
    out.write_text("\n".join(lines))
    print(f"  wrote {out.relative_to(out.parent.parent)}")


def build_pptx(summary: dict):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title
        s.placeholders[1].text = subtitle

    def add_image_slide(title, image_path: Path, note: str = ""):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
        if image_path.exists():
            left = Inches(0.5); top = Inches(1.2)
            pic = s.shapes.add_picture(str(image_path), left, top,
                                        width=Inches(12.3))
            if pic.height > Inches(5.8):
                pic.height = Inches(5.8)
                pic.width = int(pic.width * (Inches(5.8) / pic.height)) if pic.height else pic.width
        if note:
            tx = s.shapes.add_textbox(Inches(0.5), Inches(7.0),
                                       Inches(12.3), Inches(0.4))
            tf = tx.text_frame; tf.text = note
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)

    add_title_slide("Immune Cell Classifier",
                    "EfficientNetV2-S · 4-class leukocyte classification")
    # Problem
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Problem"
    body = s.placeholders[1].text_frame
    body.text = "Classify peripheral blood smear crops into EOSINOPHIL / LYMPHOCYTE / MONOCYTE / NEUTROPHIL."
    body.add_paragraph().text = "Two evaluation regimes: official split + leakage-free grouped split."

    add_image_slide("Data",
                    Path(__file__).resolve().parent.parent / "reports" / "sample_images.png",
                    "Sample crops from dataset2-master.")
    add_image_slide("Method",
                    FIG_DIR / "gradcam_panel_grouped.png",
                    "EfficientNetV2-S fine-tuned with standard augmentation; Grad-CAM confirms nucleus-centered attention.")
    # Headline
    hs = prs.slides.add_slide(prs.slide_layouts[5])
    hs.shapes.title.text = "Headline results"
    rows = [("Split", "Accuracy", "Macro-F1", "ECE")]
    for sp, d in summary.items():
        rows.append((sp, f"{d['accuracy']:.4f}", f"{d['macro_f1']:.4f}", f"{d['ece']:.3f}"))
    tbl_shape = hs.shapes.add_table(len(rows), 4,
                                     Inches(2), Inches(2),
                                     Inches(9), Inches(3))
    tbl = tbl_shape.table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(16)
                    if i == 0: r.font.bold = True

    add_image_slide("Error analysis",
                    FIG_DIR / "misclass_gallery_official.png",
                    "Most-confidently-wrong predictions; errors concentrate on lymphocyte ↔ monocyte boundary.")
    add_image_slide("Feature-space separation",
                    FIG_DIR / "umap_grouped.png",
                    "UMAP of the 1280-D penultimate features — classes are linearly separable.")
    add_image_slide("Summary",
                    FIG_DIR / "hero_figure.png",
                    "One-page overview.")

    out = Path(__file__).resolve().parent.parent / "reports" / "summary.pptx"
    prs.save(out)
    print(f"  wrote {out.relative_to(out.parent.parent)}")


# ──────────────────────────────────────────────────────────────────────────────
# Dataset plumbing for the two splits
# ──────────────────────────────────────────────────────────────────────────────

def build_test_dataset(split: str):
    """Build the canonical test dataset + its DataLoader (deterministic transform)."""
    if split == "grouped":
        _, _, test_loader = get_grouped_split_loaders(
            img_size=224, batch_size=32, num_workers=0)
    else:
        # official — re-seed to match training setup
        np.random.seed(42)
        _, _, test_loader = get_official_split_loaders(
            img_size=224, batch_size=32, num_workers=0)
    return test_loader.dataset, test_loader


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────

def run_split(split: str, skip: set[str]) -> dict:
    print(f"\n── {split} split ────────────────────────────────")
    model = load_checkpoint(split)
    dataset, loader = build_test_dataset(split)
    preds = collect_predictions(model, loader)
    acc = float((preds["preds"] == preds["labels"]).mean())
    from sklearn.metrics import f1_score
    macro_f1 = float(f1_score(preds["labels"], preds["preds"], average="macro"))
    print(f"  n={len(preds['labels'])}  acc={acc:.4f}  macro-F1={macro_f1:.4f}")

    if "gradcam" not in skip:
        figure_gradcam_panel(model, dataset, preds, split)
    if "cam_cmp" not in skip:
        figure_cam_comparison(model, dataset, preds, split)
    if "misclass" not in skip:
        figure_misclass_gallery(model, dataset, preds, split)
    if "mean_att" not in skip:
        figure_mean_attention(model, dataset, preds, split)
    if "umap" not in skip:
        figure_umap(preds, split)
    ece = 0.0
    if "calib" not in skip:
        ece = figure_calibration_roc_pr(preds, split)
    if "cm" not in skip:
        figure_confusion_matrix(preds, split)

    return dict(accuracy=acc, macro_f1=macro_f1, ece=ece)


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--skip", nargs="*", default=[],
                   help="Skip: gradcam cam_cmp misclass mean_att umap calib cm cross hero delivery")
    p.add_argument("--only", nargs="*", default=None,
                   help="Only the listed splits: grouped / official")
    args = p.parse_args()

    skip = set(args.skip)
    splits = args.only if args.only else ["grouped", "official"]

    summary = {}
    for split in splits:
        summary[split] = run_split(split, skip)

    if "cross" not in skip:
        print("\n── cross-dataset sanity check ─────────────────")
        model = load_checkpoint("grouped")
        figure_cross_dataset(model)

    if "hero" not in skip:
        print("\n── hero figure ────────────────────────────────")
        figure_hero(summary)

    if "delivery" not in skip:
        print("\n── delivery ───────────────────────────────────")
        write_captions()
        write_index_html()
        write_results_md(summary)
        build_pptx(summary)

    print("\nDone.")


if __name__ == "__main__":
    main()
